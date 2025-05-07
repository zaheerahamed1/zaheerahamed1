
import os
import re
import zipfile
import pytesseract
import tempfile
import logging
import pdfplumber
import pandas as pd
from PIL import Image
from io import BytesIO
from pptx import Presentation
from docx import Document
from win32com.client import Dispatch

# Setup logging
logging.basicConfig(filename='email_extraction.log', level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Set Tesseract path if needed
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Regex patterns
patterns = {
    "claim_number": r"(?:Claim Number|Claim #|Claim No)\s*[:\-]?\s*(\w+)",
    "npi_id": r"\bNPI\s*[:\-]?\s*(\d{10})\b",
    "recovery_item": r"Recovery Item\s*[:\-]?\s*(.+)",
    "case_id": r"(?:Case ID|Case Number)\s*[:\-]?\s*(\w+)",
    "provider_name": r"Provider Name\s*[:\-]?\s*(.+)",
    "provider_state": r"\bState\s*[:\-]?\s*([A-Z]{2})\b",
    "tax_id": r"\b(?:Tax ID|TIN)\s*[:\-]?\s*(\d{2}-\d{7}|\d{9})\b",
    "reference_number": r"(?:Ref(?:erence)? Number|Ref#)\s*[:\-]?\s*(\w+)",
}

def extract_with_patterns(text):
    result = {}
    for key, pattern in patterns.items():
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            result[key] = match.group(1).strip()
    return result

def extract_text_from_attachment(attachment):
    filename = attachment.FileName
    with tempfile.NamedTemporaryFile(delete=False) as temp:
        temp.write(attachment.Content)
        temp_path = temp.name

    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == '.txt':
            with open(temp_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == '.docx':
            doc = Document(temp_path)
            return '\n'.join(p.text for p in doc.paragraphs)
        elif ext == '.pdf':
            with pdfplumber.open(temp_path) as pdf:
                return '\n'.join(page.extract_text() for page in pdf.pages if page.extract_text())
        elif ext in ['.xlsx', '.xls']:
            df = pd.read_excel(temp_path)
            return df.to_string()
        elif ext == '.pptx':
            prs = Presentation(temp_path)
            return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
        elif ext == '.zip':
            text = ''
            with zipfile.ZipFile(temp_path, 'r') as zip_ref:
                for name in zip_ref.namelist():
                    with zip_ref.open(name) as f:
                        try:
                            content = f.read().decode('utf-8', errors='ignore')
                            text += content + "\n"
                        except:
                            continue
            return text
        else:
            return ''
    finally:
        os.remove(temp_path)

def extract_from_outlook():
    outlook = Dispatch("Outlook.Application").GetNamespace("MAPI")
    inbox = outlook.GetDefaultFolder(6)  # Inbox
    messages = inbox.Items
    extracted_data = []

    for msg in messages:
        try:
            logging.info(f"Processing email: {msg.Subject}")
            text = msg.Subject + "\n" + msg.Body

            # OCR from images
            for att in msg.Attachments:
                if att.FileName.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
                    image = Image.open(BytesIO(att.Content))
                    text += "\n" + pytesseract.image_to_string(image)

            # Other attachments
            for att in msg.Attachments:
                if not att.FileName.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
                    text += "\n" + extract_text_from_attachment(att)

            fields = extract_with_patterns(text)
            fields['email_subject'] = msg.Subject
            fields['email_sender'] = msg.SenderEmailAddress
            extracted_data.append(fields)

        except Exception as e:
            logging.error(f"Error processing email '{msg.Subject}': {e}")
            continue

    return extracted_data

if __name__ == "__main__":
    logging.info("Email extraction started.")
    results = extract_from_outlook()
    df = pd.DataFrame(results)
    output_file = "outlook_email_extracted_data.xlsx"
    df.to_excel(output_file, index=False)
    logging.info(f"Email extraction completed. Data saved to {output_file}")
    print(f"Extraction complete. Results saved to {output_file}")
